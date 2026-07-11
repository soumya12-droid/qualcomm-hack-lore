import fitz
import openpyxl
import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from pc.indexer import extractor


def test_extract_pdf_returns_one_segment_per_page(tmp_path):
    path = tmp_path / "sample.pdf"
    doc = fitz.open()
    page1 = doc.new_page()
    page1.insert_text((72, 72), "Hello PDF page one")
    page2 = doc.new_page()
    page2.insert_text((72, 72), "Hello PDF page two")
    doc.save(path)
    doc.close()

    segments = extractor.extract_pdf(path)

    assert len(segments) == 2
    assert "Hello PDF page one" in segments[0]["text"]
    assert segments[0]["page"] == "1"
    assert "Hello PDF page two" in segments[1]["text"]
    assert segments[1]["page"] == "2"
    assert segments[0]["slide"] is None and segments[0]["sheet"] is None


def test_extract_pptx_returns_one_segment_per_slide(tmp_path):
    path = tmp_path / "sample.pptx"
    presentation = Presentation()
    layout = presentation.slide_layouts[1]

    slide1 = presentation.slides.add_slide(layout)
    slide1.shapes.title.text = "First Slide"
    slide1.placeholders[1].text = "First slide body"

    slide2 = presentation.slides.add_slide(layout)
    slide2.shapes.title.text = "Second Slide"
    slide2.placeholders[1].text = "Second slide body"

    presentation.save(path)

    segments = extractor.extract_pptx(path)

    assert len(segments) == 2
    assert "First Slide" in segments[0]["text"]
    assert "First slide body" in segments[0]["text"]
    assert segments[0]["slide"] == "1"
    assert segments[1]["slide"] == "2"


def test_extract_docx_groups_paragraphs_by_heading(tmp_path):
    path = tmp_path / "sample.docx"
    document = DocxDocument()
    document.add_paragraph("Preamble before any heading")
    document.add_heading("Intro", level=1)
    document.add_paragraph("Intro body text")
    document.add_heading("Section Two", level=1)
    document.add_paragraph("Section two body text")
    document.save(path)

    segments = extractor.extract_docx(path)

    assert len(segments) == 3
    assert segments[0]["section"] is None
    assert "Preamble" in segments[0]["text"]
    assert segments[1]["section"] == "Intro"
    assert "Intro body text" in segments[1]["text"]
    assert segments[2]["section"] == "Section Two"
    assert "Section two body text" in segments[2]["text"]


def test_extract_xlsx_returns_one_segment_per_sheet(tmp_path):
    path = tmp_path / "sample.xlsx"
    workbook = openpyxl.Workbook()
    ws1 = workbook.active
    ws1.title = "Sheet1"
    ws1.append(["name", "score"])
    ws1.append(["alice", 42])

    ws2 = workbook.create_sheet("Sheet2")
    ws2.append(["other", "data"])

    workbook.save(path)

    segments = extractor.extract_xlsx(path)

    assert len(segments) == 2
    assert segments[0]["sheet"] == "Sheet1"
    assert "alice" in segments[0]["text"]
    assert "42" in segments[0]["text"]
    assert segments[1]["sheet"] == "Sheet2"
    assert "other" in segments[1]["text"]


def test_extract_txt_returns_single_segment(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text("plain text content", encoding="utf-8")

    segments = extractor.extract_text_file(path)

    assert len(segments) == 1
    assert segments[0]["text"] == "plain text content"
    assert segments[0]["page"] is None
    assert segments[0]["section"] is None


def test_extract_md_returns_single_raw_segment(tmp_path):
    path = tmp_path / "sample.md"
    path.write_text("# Title\n\nSome **bold** text.", encoding="utf-8")

    segments = extractor.extract_text_file(path)

    assert len(segments) == 1
    assert segments[0]["text"] == "# Title\n\nSome **bold** text."


def test_extract_empty_file_returns_no_segments(tmp_path):
    path = tmp_path / "empty.txt"
    path.write_text("   \n  ", encoding="utf-8")

    assert extractor.extract_text_file(path) == []


def test_strip_markdown_syntax_removes_common_markup():
    raw = "# Heading\n\nSome **bold** and *italic* text with a [link](https://example.com).\n> quoted\n- bullet"
    stripped = extractor.strip_markdown_syntax(raw)

    assert "#" not in stripped
    assert "**" not in stripped
    assert "[link]" not in stripped
    assert "(https://example.com)" not in stripped
    assert "link" in stripped
    assert "bold" in stripped
    assert "italic" in stripped
    assert "quoted" in stripped
    assert "bullet" in stripped


def test_extract_dispatches_by_extension(tmp_path):
    path = tmp_path / "sample.txt"
    path.write_text("dispatched content", encoding="utf-8")

    segments = extractor.extract(path)

    assert segments == [{"text": "dispatched content", "page": None, "slide": None, "sheet": None, "section": None}]


def test_extract_raises_for_unsupported_extension(tmp_path):
    path = tmp_path / "sample.unsupported"
    path.write_text("data", encoding="utf-8")

    with pytest.raises(extractor.UnsupportedFileTypeError):
        extractor.extract(path)


@pytest.mark.parametrize(
    "filename,expected",
    [
        ("a.pdf", "pdf"),
        ("a.pptx", "pptx"),
        ("a.ppt", "pptx"),
        ("a.docx", "docx"),
        ("a.xlsx", "xlsx"),
        ("a.md", "md"),
        ("a.txt", "txt"),
        ("a.unknown", None),
    ],
)
def test_file_type_for(filename, expected):
    assert extractor.file_type_for(filename) == expected
