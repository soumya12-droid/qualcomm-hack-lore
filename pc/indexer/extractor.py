"""Phase 1 — extracts text and structural context from supported file formats.

Supported formats: .pdf (PyMuPDF), .pptx/.ppt (python-pptx), .docx
(python-docx), .xlsx (openpyxl), .md and .txt (built-in).

Input: a file path.
Output: `extract()` returns a list of segment dicts, one per structural unit
of the source document (PDF page, PPTX slide, XLSX sheet, DOCX/MD section):
    {"text": str, "page": str | None, "slide": str | None,
     "sheet": str | None, "section": str | None}
Side effects: reads the file from disk.
"""

from pathlib import Path
import re

import fitz  # PyMuPDF
import openpyxl
from docx import Document as DocxDocument
from pptx import Presentation

FILE_TYPE_BY_EXTENSION = {
    ".pdf": "pdf",
    ".pptx": "pptx",
    ".ppt": "pptx",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".md": "md",
    ".txt": "txt",
}


class UnsupportedFileTypeError(ValueError):
    """Raised by extract() when the file extension has no registered extractor."""


def _segment(text, page=None, slide=None, sheet=None, section=None):
    return {"text": text, "page": page, "slide": slide, "sheet": sheet, "section": section}


def extract_pdf(file_path):
    """Extract one segment per non-empty page. Side effect: opens the PDF via PyMuPDF."""
    segments = []
    with fitz.open(file_path) as doc:
        for index, page in enumerate(doc):
            text = page.get_text().strip()
            if text:
                segments.append(_segment(text, page=str(index + 1)))
    return segments


def extract_pptx(file_path):
    """Extract one segment per non-empty slide, concatenating all shape text.

    Note: python-pptx only reads the OOXML .pptx format. Legacy binary .ppt
    files are registered under the same file_type per CLAUDE.md but will
    raise when opened here — convert legacy .ppt to .pptx before indexing.
    """
    segments = []
    presentation = Presentation(file_path)
    for index, slide in enumerate(presentation.slides):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
        text = "\n".join(parts).strip()
        if text:
            segments.append(_segment(text, slide=str(index + 1)))
    return segments


def extract_docx(file_path):
    """Extract one segment per section, where a section is the run of body
    paragraphs following a "Heading*"-styled paragraph (or the document start
    if no heading precedes it)."""
    document = DocxDocument(file_path)
    segments = []
    current_section = None
    buffer = []

    def flush():
        text = "\n".join(buffer).strip()
        if text:
            segments.append(_segment(text, section=current_section))

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        is_heading = paragraph.style is not None and paragraph.style.name.startswith("Heading")
        if is_heading:
            flush()
            buffer = []
            current_section = text
        else:
            buffer.append(text)
    flush()
    return segments


def extract_xlsx(file_path):
    """Extract one segment per non-empty sheet, rendering rows as tab-separated lines."""
    segments = []
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    for sheet in workbook.worksheets:
        lines = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                lines.append("\t".join(cells))
        text = "\n".join(lines).strip()
        if text:
            segments.append(_segment(text, sheet=sheet.title))
    return segments


def extract_text_file(file_path):
    """Extract .md/.txt files as a single raw-text segment (no structural subdivision)."""
    text = Path(file_path).read_text(encoding="utf-8").strip()
    if not text:
        return []
    return [_segment(text)]


_EXTRACTORS = {
    ".pdf": extract_pdf,
    ".pptx": extract_pptx,
    ".ppt": extract_pptx,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".md": extract_text_file,
    ".txt": extract_text_file,
}

_MARKDOWN_PATTERNS = [
    (re.compile(r"^#{1,6}\s+", re.MULTILINE), ""),        # headings
    (re.compile(r"`{1,3}([^`]+)`{1,3}"), r"\1"),           # inline code / fences
    (re.compile(r"\*\*([^*]+)\*\*"), r"\1"),               # bold
    (re.compile(r"\*([^*]+)\*"), r"\1"),                   # italics
    (re.compile(r"\[([^\]]+)\]\([^)]+\)"), r"\1"),         # links
    (re.compile(r"^>\s?", re.MULTILINE), ""),              # blockquotes
    (re.compile(r"^[-*+]\s+", re.MULTILINE), ""),          # bullet list markers
]


def strip_markdown_syntax(text):
    """Strip common Markdown syntax, for use as embedding-time text on .md
    chunks (the raw chunk text with syntax intact is still what gets stored,
    per CLAUDE.md's "strip for embedding, keep raw for storage")."""
    for pattern, replacement in _MARKDOWN_PATTERNS:
        text = pattern.sub(replacement, text)
    return text


def file_type_for(file_path):
    """Return the LanceDB `file_type` value for a path's extension, or None if unsupported."""
    return FILE_TYPE_BY_EXTENSION.get(Path(file_path).suffix.lower())


def extract(file_path):
    """Dispatch to the extractor registered for `file_path`'s extension.

    Raises:
        UnsupportedFileTypeError: if the extension has no registered extractor.
    """
    path = Path(file_path)
    extractor_fn = _EXTRACTORS.get(path.suffix.lower())
    if extractor_fn is None:
        raise UnsupportedFileTypeError(f"Unsupported file extension: {path.suffix!r}")
    return extractor_fn(path)
